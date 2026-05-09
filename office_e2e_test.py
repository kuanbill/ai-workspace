import json
import os
import sys
import tempfile
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent
VENDOR_DIR = BASE_DIR / "vendor"
if VENDOR_DIR.exists():
    sys.path.insert(0, str(VENDOR_DIR))

from docx import Document as DocxDocument
from openpyxl import Workbook
from pypdf import PdfWriter
from pptx import Presentation

from tools import handle_tool_call


def call_tool(project_root: str, name: str, args: dict) -> dict:
    raw = handle_tool_call(name, args, project_root)
    try:
        return json.loads(raw)
    except Exception:
        return {"raw": raw}


def create_samples(project_dir: Path) -> None:
    doc = DocxDocument()
    doc.add_heading("Office E2E", 0)
    doc.add_paragraph("This is a DOCX test file.")
    doc.save(project_dir / "sample.docx")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Office E2E"
    slide.placeholders[1].text = "PPTX test file"
    prs.save(project_dir / "sample.pptx")

    pdf = PdfWriter()
    pdf.add_blank_page(width=300, height=200)
    with open(project_dir / "sample.pdf", "wb") as f:
        pdf.write(f)

    wb = Workbook()
    ws = wb.active
    ws["A1"] = 1
    ws["A2"] = 2
    ws["A3"] = "=SUM(A1:A2)"
    wb.save(project_dir / "sample.xlsx")

    (project_dir / "hello.py").write_text("print('python-ok')\n", encoding="utf-8")
    (project_dir / "hello.js").write_text("console.log('node-ok')\n", encoding="utf-8")
    (project_dir / "pdf_empty_fields.json").write_text("[]", encoding="utf-8")


def main() -> None:
    project_root = Path(tempfile.mkdtemp(prefix="office_e2e_"))
    create_samples(project_root)

    results = {
        "project_root": str(project_root),
        "file_tools": {},
        "office_tools": {},
    }

    results["file_tools"]["write_file"] = call_tool(
        str(project_root), "write_file", {"path": "notes.txt", "content": "abc"}
    )
    results["file_tools"]["read_file"] = call_tool(
        str(project_root), "read_file", {"path": "notes.txt"}
    )
    results["file_tools"]["list_files"] = call_tool(
        str(project_root), "list_files", {"path": ""}
    )
    results["file_tools"]["run_project_python"] = call_tool(
        str(project_root), "run_project_python", {"path": "hello.py"}
    )
    results["file_tools"]["run_project_node"] = call_tool(
        str(project_root), "run_project_node", {"path": "hello.js"}
    )

    results["office_tools"]["env"] = call_tool(
        str(project_root), "get_office_environment_status", {}
    )
    results["office_tools"]["skill_doc"] = call_tool(
        str(project_root), "read_office_skill_doc", {"doc_id": "pptx_skill"}
    )

    results["office_tools"]["docx_unpack"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "docx_unpack", "paths": {"input_path": "sample.docx", "output_dir": "docx_unpacked"}},
    )
    results["office_tools"]["docx_validate"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "docx_validate", "paths": {"input_dir": "docx_unpacked", "input_path": "sample.docx"}},
    )
    results["office_tools"]["docx_pack"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "docx_pack", "paths": {"input_dir": "docx_unpacked", "output_path": "sample_repacked.docx"}},
    )

    results["office_tools"]["pptx_unpack"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "pptx_unpack", "paths": {"input_path": "sample.pptx", "output_dir": "pptx_unpacked"}},
    )
    results["office_tools"]["pptx_validate"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "pptx_validate", "paths": {"input_dir": "pptx_unpacked", "input_path": "sample.pptx"}},
    )
    results["office_tools"]["pptx_pack"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "pptx_pack", "paths": {"input_dir": "pptx_unpacked", "output_path": "sample_repacked.pptx"}},
    )
    results["office_tools"]["pptx_inventory"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "pptx_inventory", "paths": {"input_path": "sample.pptx", "output_path": "pptx_inventory.json"}},
    )
    results["office_tools"]["pptx_thumbnail"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "pptx_thumbnail", "paths": {"input_path": "sample.pptx", "output_path": "pptx_thumbs"}},
    )

    results["office_tools"]["pdf_extract_form_fields"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "pdf_extract_form_fields", "paths": {"input_path": "sample.pdf", "output_path": "pdf_fields.json"}},
    )
    results["office_tools"]["pdf_fill_form_fields"] = call_tool(
        str(project_root),
        "run_office_script",
        {
            "script_id": "pdf_fill_form_fields",
            "paths": {
                "input_path": "sample.pdf",
                "json_path": "pdf_empty_fields.json",
                "output_path": "sample_filled.pdf",
            },
        },
    )
    results["office_tools"]["pdf_convert_to_images"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "pdf_convert_to_images", "paths": {"input_path": "sample.pdf", "output_dir": "pdf_images"}},
    )

    results["office_tools"]["xlsx_recalc"] = call_tool(
        str(project_root),
        "run_office_script",
        {"script_id": "xlsx_recalc", "paths": {"input_path": "sample.xlsx"}},
    )

    print(json.dumps(results, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
